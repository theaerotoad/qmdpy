"""
Document and Text Converters for QMD.

Converts non-markdown formats (.docx, .pptx, .xlsx, .csv, .html) into Markdown text.
"""
import csv
import io
import re
from collections import Counter
from pathlib import Path
from typing import Union, List, Optional, Dict, Any

def _get_dplib():
    try:
        import dplib
        return dplib
    except ImportError:
        import sys
        src_dir = str(Path(__file__).resolve().parent.parent)
        if src_dir not in sys.path:
            sys.path.insert(0, src_dir)
        try:
            import dplib
            return dplib
        except ImportError:
            return None

SUPPORTED_EXTENSIONS = {
    ".md", ".markdown", ".txt",
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".csv",
    ".html", ".htm",
    ".epub"
}


def _mathml_to_latex(node) -> str:
    try:
        import mathml_to_latex
        from lxml import etree
        math_xml = etree.tostring(node, encoding='unicode')
        if hasattr(mathml_to_latex, 'convert'):
            return mathml_to_latex.convert(math_xml)
        elif hasattr(mathml_to_latex, 'mathml_to_latex'):
            return mathml_to_latex.mathml_to_latex(math_xml)
        else:
            return "".join(node.itertext())
    except ImportError:
        return "".join(node.itertext())
    except Exception:
        return "".join(node.itertext())

def _omml_to_latex(node) -> str:
    MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    res = ""
    for child in node:
        if not isinstance(child.tag, str):
            continue
        tag = child.tag.split('}')[-1]
        if tag == 'f':
            num = child.find(f"{{{MATH_NS}}}num")
            den = child.find(f"{{{MATH_NS}}}den")
            num_tex = _omml_to_latex(num) if num is not None else ""
            den_tex = _omml_to_latex(den) if den is not None else ""
            res += f"\\frac{{{num_tex}}}{{{den_tex}}}"
        elif tag == 'rad':
            deg = child.find(f"{{{MATH_NS}}}deg")
            e = child.find(f"{{{MATH_NS}}}e")
            deg_tex = _omml_to_latex(deg) if deg is not None and len(deg) else ""
            e_tex = _omml_to_latex(e) if e is not None else ""
            if deg_tex:
                res += f"\\sqrt[{deg_tex}]{{{e_tex}}}"
            else:
                res += f"\\sqrt{{{e_tex}}}"
        elif tag == 'sSup':
            e = child.find(f"{{{MATH_NS}}}e")
            sup = child.find(f"{{{MATH_NS}}}sup")
            e_tex = _omml_to_latex(e) if e is not None else ""
            sup_tex = _omml_to_latex(sup) if sup is not None else ""
            res += f"{e_tex}^{{{sup_tex}}}"
        elif tag == 'sSub':
            e = child.find(f"{{{MATH_NS}}}e")
            sub = child.find(f"{{{MATH_NS}}}sub")
            e_tex = _omml_to_latex(e) if e is not None else ""
            sub_tex = _omml_to_latex(sub) if sub is not None else ""
            res += f"{e_tex}_{{{sub_tex}}}"
        elif tag == 'sSubSup':
            e = child.find(f"{{{MATH_NS}}}e")
            sub = child.find(f"{{{MATH_NS}}}sub")
            sup = child.find(f"{{{MATH_NS}}}sup")
            e_tex = _omml_to_latex(e) if e is not None else ""
            sub_tex = _omml_to_latex(sub) if sub is not None else ""
            sup_tex = _omml_to_latex(sup) if sup is not None else ""
            res += f"{e_tex}_{{{sub_tex}}}^{{{sup_tex}}}"
        elif tag == 'd':
            dPr = child.find(f"{{{MATH_NS}}}dPr")
            beg_chr = "("
            end_chr = ")"
            sep_chr = "|"
            if dPr is not None:
                beg_node = dPr.find(f"{{{MATH_NS}}}begChr")
                if beg_node is not None: beg_chr = beg_node.get(f"{{{MATH_NS}}}val", "(")
                end_node = dPr.find(f"{{{MATH_NS}}}endChr")
                if end_node is not None: end_chr = end_node.get(f"{{{MATH_NS}}}val", ")")
                sep_node = dPr.find(f"{{{MATH_NS}}}sepChr")
                if sep_node is not None: sep_chr = sep_node.get(f"{{{MATH_NS}}}val", "|")
            
            if beg_chr == "{": beg_chr = "\\{"
            if end_chr == "}": end_chr = "\\}"
            
            e_nodes = child.findall(f"{{{MATH_NS}}}e")
            e_tex = f" {sep_chr} ".join(_omml_to_latex(e) for e in e_nodes)
            res += f"\\left{beg_chr} {e_tex} \\right{end_chr}"
        elif tag == 'nary':
            naryPr = child.find(f"{{{MATH_NS}}}naryPr")
            op_tex = "\\int"
            if naryPr is not None:
                chr_node = naryPr.find(f"{{{MATH_NS}}}chr")
                if chr_node is not None:
                    val = chr_node.get(f"{{{MATH_NS}}}val")
                    if val == "∑": op_tex = "\\sum"
                    elif val == "∏": op_tex = "\\prod"
                    elif val == "∐": op_tex = "\\coprod"
                    elif val == "∪": op_tex = "\\bigcup"
                    elif val == "∩": op_tex = "\\bigcap"

            sub = child.find(f"{{{MATH_NS}}}sub")
            sup = child.find(f"{{{MATH_NS}}}sup")
            e = child.find(f"{{{MATH_NS}}}e")
            sub_tex = _omml_to_latex(sub) if sub is not None else ""
            sup_tex = _omml_to_latex(sup) if sup is not None else ""
            e_tex = _omml_to_latex(e) if e is not None else ""
            
            res += f"{op_tex}"
            if sub_tex: res += f"_{{{sub_tex}}}"
            if sup_tex: res += f"^{{{sup_tex}}}"
            res += f" {e_tex}"
        elif tag == 'func':
            fName = child.find(f"{{{MATH_NS}}}fName")
            e = child.find(f"{{{MATH_NS}}}e")
            fName_tex = _omml_to_latex(fName) if fName is not None else ""
            e_tex = _omml_to_latex(e) if e is not None else ""
            res += f"{fName_tex}{e_tex}"
        elif tag == 'm':
            mr_nodes = child.findall(f"{{{MATH_NS}}}mr")
            rows = []
            for mr in mr_nodes:
                e_nodes = mr.findall(f"{{{MATH_NS}}}e")
                rows.append(" & ".join(_omml_to_latex(e) for e in e_nodes))
            matrix_content = " \\\\ ".join(rows)
            res += f"\\begin{{matrix}} {matrix_content} \\end{{matrix}}"
        elif tag == 'eqArr':
            e_nodes = child.findall(f"{{{MATH_NS}}}e")
            lines = [ _omml_to_latex(e) for e in e_nodes ]
            res += " \\\\ ".join(lines)
        elif tag == 'acc':
            accPr = child.find(f"{{{MATH_NS}}}accPr")
            chr_val = "^"
            if accPr is not None:
                chr_node = accPr.find(f"{{{MATH_NS}}}chr")
                if chr_node is not None:
                    chr_val = chr_node.get(f"{{{MATH_NS}}}val", "^")
            e = child.find(f"{{{MATH_NS}}}e")
            e_tex = _omml_to_latex(e) if e is not None else ""
            if chr_val in ("⃗", "\u20d7"): res += f"\\vec{{{e_tex}}}"
            elif chr_val in ("^", "̂", "\u0302"): res += f"\\hat{{{e_tex}}}"
            elif chr_val in ("‾", "¯", "\u0304"): res += f"\\bar{{{e_tex}}}"
            elif chr_val in ("˜", "̃", "\u0303"): res += f"\\tilde{{{e_tex}}}"
            elif chr_val in ("˙", "̇", "\u0307"): res += f"\\dot{{{e_tex}}}"
            elif chr_val in ("¨", "̈", "\u0308"): res += f"\\ddot{{{e_tex}}}"
            else: res += f"{e_tex}"
        elif tag == 'limLow':
            e = child.find(f"{{{MATH_NS}}}e")
            lim = child.find(f"{{{MATH_NS}}}lim")
            e_tex = _omml_to_latex(e) if e is not None else ""
            lim_tex = _omml_to_latex(lim) if lim is not None else ""
            res += f"\\mathop{{{e_tex}}}_{{{lim_tex}}}"
        elif tag == 't':
            res += child.text or ""
        else:
            res += _omml_to_latex(child)
    return res

def _extract_text_and_math(node) -> str:
    """Recursively extracts text and math formulas from an lxml node (e.g. DOCX CT_P)."""
    MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
    MML_NS = "http://www.w3.org/1998/Math/MathML"
    WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    
    if not hasattr(node, "tag"):
        return ""

    tag = node.tag
    if isinstance(tag, str):
        if tag == f"{{{MATH_NS}}}oMathPara":
            tex = _omml_to_latex(node)
            return f"$$ {tex} $$"
        elif tag == f"{{{MATH_NS}}}oMath":
            tex = _omml_to_latex(node)
            return f"${tex}$"
        elif tag == f"{{{MML_NS}}}math":
            tex = _mathml_to_latex(node)
            return f"${tex}$"
        
        # Explicit text nodes: Word, Word-Delete, Math, PPTX
        elif tag in (f"{{{WORD_NS}}}t", f"{{{WORD_NS}}}delText", f"{{{MATH_NS}}}t", f"{{{DRAWING_NS}}}t"):
            return node.text or ""
        # Handle whitespace explicitly
        elif tag == f"{{{WORD_NS}}}tab":
            return "\t"
        elif tag == f"{{{WORD_NS}}}br":
            return "\n"
    
    res = []
    for child in node:
        res.append(_extract_text_and_math(child))
        
    return "".join(res)


def is_supported_file(file_path: Union[str, Path]) -> bool:
    ext = Path(file_path).suffix.lower()
    return ext in SUPPORTED_EXTENSIONS


def _sanitize_text(text: str) -> str:
    """
    Sanitizes string by replacing invalid unicode surrogates (U+D800 - U+DFFF)
    with replacement characters so it can be encoded to UTF-8 without errors.
    """
    if not text:
        return ""
    try:
        return text.encode('utf-16', 'surrogatepass').decode('utf-16', 'replace')
    except Exception:
        return text


def _is_image_processing_enabled(config) -> bool:
    if not config:
        return False
    return bool(
        getattr(config, "vision_url", None)
        or getattr(config, "multimodal_url", None)
        or getattr(config, "multimodal_model", None)
    )


def _process_image_multimodal_llm(image_bytes: bytes, filename: str, config, errors_out: Optional[List[dict]] = None) -> str:
    if not config or not image_bytes:
        return ""
    try:
        from qmd.llm import LLMClient
        client = LLMClient(
            base_url=getattr(config, "llm_url", None),
            api_key=getattr(config, "api_key", None),
            multimodal_url=getattr(config, "multimodal_url", None),
            multimodal_api_key=getattr(config, "multimodal_api_key", None),
            multimodal_model=getattr(config, "multimodal_model", None),
            multimodal_prompt=getattr(config, "multimodal_prompt", None),
            timeout=getattr(config, "request_timeout", 120.0),
        )
        return client.process_image(image_bytes, filename=filename)
    except Exception as e:
        print(f"Warning: Multimodal LLM error for {filename}: {e}")
        if errors_out is not None:
            errors_out.append({"error_type": "multimodal_image_error", "message": f"{filename}: {e}"})
        return ""


def _process_image(image_bytes: bytes, filename: str, config, errors_out: Optional[List[dict]] = None) -> str:
    if not config or not image_bytes:
        return ""
    if getattr(config, "multimodal_url", None) or getattr(config, "multimodal_model", None):
        try:
            return _process_image_multimodal_llm(image_bytes, filename, config, errors_out=errors_out)
        except TypeError:
            return _process_image_multimodal_llm(image_bytes, filename, config)
    elif getattr(config, "vision_url", None):
        try:
            return _process_image_vision_api(image_bytes, filename, config, errors_out=errors_out)
        except TypeError:
            return _process_image_vision_api(image_bytes, filename, config)
    return ""


def _process_images_concurrently(
    items: List[tuple],
    config,
    max_workers: Union[int, None] = None,
    errors_out: Optional[List[dict]] = None
) -> List[str]:
    if not items:
        return []
    if max_workers is None:
        max_workers = (
            getattr(config, "max_image_concurrency", None)
            or getattr(config, "max_simultaneous_images", None)
            or 4
        )
    max_workers = max(1, int(max_workers))

    if len(items) == 1 or max_workers == 1:
        return [_process_image(b, fn, config, errors_out=errors_out) for b, fn in items]

    import concurrent.futures
    workers = min(len(items), max_workers)
    with concurrent.futures.ThreadPoolExecutor(max_workers=workers) as executor:
        futures = [executor.submit(_process_image, b, fn, config, errors_out) for b, fn in items]
        results = []
        for f in futures:
            try:
                results.append(f.result())
            except Exception as e:
                print(f"Warning: Concurrent image processing error: {e}")
                if errors_out is not None:
                    errors_out.append({"error_type": "image_processing_error", "message": str(e)})
                results.append("")
        return results


def _process_image_vision_api(image_bytes: bytes, filename: str, config, errors_out: Optional[List[dict]] = None) -> str:
    if not config or not getattr(config, "vision_url", None):
        return ""
    
    try:
        import httpx
        import base64
    except ImportError:
        return ""
    
    vision_url = config.vision_url
    headers = {}
    if getattr(config, "vision_api_key", None):
        headers["Authorization"] = f"Bearer {config.vision_api_key}"
        
    b64_image = base64.b64encode(image_bytes).decode('utf-8')
    mime_type = "image/jpeg"
    if filename.lower().endswith(".png"): mime_type = "image/png"
    
    payload = {
        "image": f"data:{mime_type};base64,{b64_image}",
        "extract_tables": True,
        "extract_text": True
    }
    
    try:
        resp = httpx.post(vision_url, json=payload, headers=headers, timeout=getattr(config, "request_timeout", 120.0))
        resp.raise_for_status()
        data = resp.json()
        
        detections = data.get("detections", [])
        if not detections and "results" in data and len(data["results"]) > 0:
            detections = data["results"][0].get("detections", [])
            
        md_lines = []
        text_labels = {
            "caption", "footnote", "formula", "list-item", 
            "page-footer", "page-header", "section-header", 
            "text", "title"
        }
        
        for d in detections:
            lbl = d.get("label", "").lower()
            if lbl == "picture":
                text = d.get("text", "").strip()
                alt = f"Image with text: {text}" if text else "Image"
                md_lines.append(f"![{alt}]({filename})")
            elif lbl == "table":
                if d.get("markdown"):
                    md_lines.append(d["markdown"])
                elif d.get("html"):
                    md_lines.append(d["html"])
            elif lbl in text_labels:
                if d.get("text"):
                    md_lines.append(d["text"].strip())
                    
        return "\n\n".join(md_lines)
    except Exception as e:
        print(f"Warning: Vision API error for {filename}: {e}")
        if errors_out is not None:
            errors_out.append({"error_type": "vision_api_error", "message": f"{filename}: {e}"})
        return ""


def convert_to_markdown(file_path: Union[str, Path], config=None, errors_out: Optional[List[dict]] = None) -> str:
    """
    Converts a supported document or text file to Markdown.
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext in {".md", ".markdown", ".txt"}:
        raw_md = _convert_text(path)
    elif ext == ".pdf":
        raw_md = _convert_pdf(path, config, errors_out=errors_out)
    elif ext == ".docx":
        raw_md = _convert_docx(path, config, errors_out=errors_out)
    elif ext == ".pptx":
        raw_md = _convert_pptx(path, config, errors_out=errors_out)
    elif ext == ".xlsx":
        raw_md = _convert_xlsx(path, config, errors_out=errors_out)
    elif ext == ".csv":
        raw_md = _convert_csv(path)
    elif ext in {".html", ".htm"}:
        raw_md = _convert_html(path)
    elif ext == ".epub":
        raw_md = _convert_epub(path)
    else:
        raw_md = _convert_text(path)

    sanitized = _sanitize_text(raw_md)
    
    # Global pass: Condense extra padding spaces in any markdown tables to save tokens
    # Safely target only contiguous blocks containing a markdown table separator (|---|)
    def _condense_table_block(match):
        block = match.group(0)
        # Check if block contains a valid markdown separator line, allowing leading/trailing spaces
        if re.search(r'^[ \t]*\|[-\s:|]+\|[ \t]*$', block, flags=re.MULTILINE):
            # Replace 2 or more spaces with a single space
            return re.sub(r' {2,}', ' ', block)
        return block
        
    return re.sub(r'(?:^[ \t]*\|[^\n]*\|[ \t]*(?:\r?\n|$))+', _condense_table_block, sanitized, flags=re.MULTILINE)


def _convert_text(path: Path) -> str:
    raw_bytes = path.read_bytes()

    # Reject binary files containing null bytes in header/first chunk
    if b'\x00' in raw_bytes[:8192]:
        raise ValueError(f"File '{path.name}' appears to be a binary file and cannot be converted as plain text.")

    for enc in ("utf-8", "utf-8-sig"):
        try:
            return raw_bytes.decode(enc)
        except UnicodeDecodeError:
            continue

    try:
        return raw_bytes.decode("latin-1")
    except Exception as e:
        raise ValueError(f"Could not decode text file '{path.name}': {e}")


def _convert_pdf(path: Path, config=None, errors_out: Optional[List[dict]] = None) -> str:
    import os, sys
    verbose = os.environ.get("QMD_VERBOSE") == "1"
    
    try:
        import pymupdf
        import pymupdf4llm
    except ImportError:
        raise ImportError("pymupdf and pymupdf4llm are required for converting .pdf files. Install with `pip install pymupdf pymupdf4llm`.")

    # Silence PyMuPDF C-level warnings globally before doing anything
    if hasattr(pymupdf, "TOOLS"):
        pymupdf.TOOLS.mupdf_display_errors(False)

    if verbose:
        print(f"[Verbose PDF] Attempting to open {path} with pymupdf...", flush=True)

    doc = pymupdf.open(str(path))
    
    # 0. Sanitize the PDF to fix corrupted xrefs/colorspaces before pymupdf4llm chokes
    try:
        # garbage=4 removes unreferenced objects, clean=True sanitizes content streams 
        # (fixes "Line cannot be recognized" and many colorspace issues)
        sanitized_bytes = doc.tobytes(garbage=4, clean=True, deflate=True)
        doc.close()
        doc = pymupdf.open(stream=sanitized_bytes, filetype="pdf")
        if verbose:
            print(f"[Verbose PDF] Successfully sanitized {path} in memory.", flush=True)
    except Exception as e:
        if verbose:
            print(f"[Verbose PDF] Pre-sanitization failed ({e}), proceeding with original.", flush=True)
        doc = pymupdf.open(str(path))

    if verbose:
        print(f"[Verbose PDF] Successfully opened {path}. Parsing headers...", flush=True)
        
    # 1. Build header detector
    def build_header_detector(doc, max_levels=5):
        toc = doc.get_toc()
        if toc:
            toc_by_page = {}
            for lvl, title, page_num in toc:
                p_idx = page_num - 1
                cleaned_title = title.strip().lower()
                if cleaned_title:
                    toc_by_page.setdefault(p_idx, []).append((min(lvl, max_levels), cleaned_title))

            def toc_header_fn(span, page=None):
                if page is None: return ""
                page_num = page.number
                if page_num not in toc_by_page: return ""
                text = span.get("text", "").strip().lower()
                if len(text) <= 2: return ""
                for lvl, title in toc_by_page[page_num]:
                    if text == title or text.startswith(title) or title.startswith(text):
                        return "#" * lvl + " "
                return ""
            return toc_header_fn

        font_sizes = Counter()
        for page in doc:
            page_dict = page.get_text("dict")
            for block in page_dict.get("blocks", []):
                if block.get("type") == 0:
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            text = span.get("text", "").strip()
                            if len(text) > 2:
                                size = round(span.get("size", 0))
                                font_sizes[size] += len(text)
        if not font_sizes:
            return None

        body_size = font_sizes.most_common(1)[0][0]
        larger_sizes = sorted([s for s in font_sizes.keys() if s > body_size], reverse=True)
        header_mapping = {size: min(idx + 1, max_levels) for idx, size in enumerate(larger_sizes[:max_levels])}

        def font_size_header_fn(span, page=None):
            text = span.get("text", "").strip()
            if len(text) <= 2: return ""
            size = round(span.get("size", 0))
            if size in header_mapping:
                return "#" * header_mapping[size] + " "
            return ""

        return font_size_header_fn

    hdr_fn = build_header_detector(doc, max_levels=5)

    # 2. Extract markdown using pymupdf4llm (suppressing noisy OCR stdout)
    import contextlib
    import io
    
    if verbose:
        print(f"[Verbose PDF] Executing pymupdf4llm.to_markdown on {path}...", flush=True)

    f = io.StringIO()
    try:
        with contextlib.redirect_stdout(f), contextlib.redirect_stderr(f):
            try:
                page_chunks = pymupdf4llm.to_markdown(
                    doc,
                    hdr_info=hdr_fn,
                    header=False,
                    footer=False,
                    write_images=False,
                    page_chunks=True,
                    show_progress=False
                )
            except TypeError:
                # Fallback for older versions of pymupdf4llm that lack show_progress
                page_chunks = pymupdf4llm.to_markdown(
                    doc,
                    hdr_info=hdr_fn,
                    header=False,
                    footer=False,
                    write_images=False,
                    page_chunks=True
                )
        if verbose:
            print(f"[Verbose PDF] Successfully parsed PDF {path} to markdown.", flush=True)
    except Exception as e:
        if verbose:
            print(f"[Verbose PDF] pymupdf4llm failed ({e}). Falling back to layout-aware text extraction for {path}...", flush=True)
        
        page_chunks = []
        for i in range(len(doc)):
            page = doc[i]
            page_md = []
            try:
                # Use native dict extraction to ignore broken images but preserve block grouping
                page_dict = page.get_text("dict")
                for block in page_dict.get("blocks", []):
                    if block.get("type") == 0:  # text block
                        block_lines = []
                        header_prefix = ""
                        for line in block.get("lines", []):
                            line_text = ""
                            for span in line.get("spans", []):
                                text = span.get("text", "")
                                if text.strip():
                                    # Reuse our header detector to reconstruct Markdown structure
                                    prefix = hdr_fn(span, page)
                                    if prefix and not header_prefix:
                                        header_prefix = prefix
                                    line_text += text
                            if line_text.strip():
                                block_lines.append(line_text.strip())
                        
                        if block_lines:
                            merged = " ".join(block_lines)
                            if header_prefix:
                                page_md.append(f"\n{header_prefix}{merged}\n")
                            else:
                                page_md.append(f"{merged}\n")
                
                final_page_text = "\n".join(page_md) if page_md else page.get_text("text")
            except Exception:
                # Ultimate fallback if dict extraction somehow fails
                final_page_text = page.get_text("text")
                
            page_chunks.append({"text": final_page_text})

        # Only record an error if layout-aware fallback also failed to extract any text
        if errors_out is not None and not any(chunk.get("text", "").strip() for chunk in page_chunks):
            errors_out.append({
                "error_type": "pdf_extraction_failed",
                "message": f"{path.name}: pymupdf4llm failed ({e}) and layout fallback extraction yielded no text"
            })
    
    seen_xrefs = set()
    
    if isinstance(page_chunks, str):
        page_chunks = [{"text": page_chunks}]

    md_pages = []
    for i, chunk in enumerate(page_chunks):
        page_text = chunk.get("text", "")
        
        # Convert pymupdf4llm's native picture text blocks to standard markdown image tags
        def _format_pic_text(match):
            text = match.group(1).strip()
            text = re.sub(r'\s+', ' ', text)
            alt = f"Image with text: {text}" if text else "Image"
            return f"\n![{alt}](pdf_image.png)\n"
            
        page_text = re.sub(r'<!--\s*Start of picture text\s*-->(.*?)<!--\s*End of picture text\s*-->\n*', _format_pic_text, page_text, flags=re.DOTALL | re.IGNORECASE)
        
        # Clean up any remaining HTML comments from pymupdf4llm
        page_text = re.sub(r'<!--.*?-->\n*', '', page_text, flags=re.DOTALL)
        
        md_pages.append(page_text)

    # Process all embedded PDF images concurrently across pages if configured
    if config and _is_image_processing_enabled(config):
        try:
            pdf_images = []
            for i in range(len(doc)):
                page = doc[i]
                for img in page.get_images():
                    xref = img[0]
                    if xref in seen_xrefs:
                        continue
                    seen_xrefs.add(xref)
                    try:
                        base_image = doc.extract_image(xref)
                    except Exception:
                        continue
                    if not base_image:
                        continue
                    image_bytes = base_image.get("image")
                    if not image_bytes:
                        continue
                    width = base_image.get("width", 0)
                    height = base_image.get("height", 0)
                    # Skip spacer/divider lines (<= 2px) that trigger MuPDF colorspace or scale errors
                    if width > 0 and height > 0 and (width <= 2 or height <= 2):
                        continue
                    ext = base_image.get("ext", "png")
                    filename = f"page_{i+1}_img_{xref}.{ext}"
                    pdf_images.append((i, image_bytes, filename))

            if pdf_images:
                image_inputs = [(img_bytes, fn) for _, img_bytes, fn in pdf_images]
                results = _process_images_concurrently(image_inputs, config, errors_out=errors_out)
                for (page_idx, _, _), img_md in zip(pdf_images, results):
                    if img_md and page_idx < len(md_pages):
                        md_pages[page_idx] += f"\n\n{img_md}\n"
        except Exception:
            pass
        
    raw_md = "\n\n".join(md_pages)

    doc.close()

    # 3. Clean and normalize headings
    def clean_heading_text(raw_heading: str) -> str | None:
        text = re.sub(r'<[^>]+>', '', raw_heading)
        text = re.sub(r'[*_~`]', '', text)
        text = re.sub(r'\s+', ' ', text).strip()
        text = text.strip("•·-–— \t")
        if not re.search(r'[a-zA-Z0-9]', text):
            return None
        return text

    def normalize_markdown_headings(markdown_text: str, book_title: str | None = None) -> str:
        text = re.sub(r'^(?:#{1,6})\s+(\*\*(?:FIGURE|TABLE|EXHIBIT|CHART|STEP\s+\d+).*?\*\*)\s*$', r'\1', markdown_text, flags=re.MULTILINE | re.IGNORECASE)
        text = re.sub(r'^#{1,6}\s+([A-Z])\s*$', r'\1', text, flags=re.MULTILINE)
        text = re.sub(r'^(?:#{1,6}\s+)?Chapter\s+(?:<u>)?([0-9]+|[IVXLCDM]+)(?:</u>)?\s*\n+(?:#{1,6}\s+)?([^\n]+)', r'## Chapter \1: \2', text, flags=re.MULTILINE | re.IGNORECASE)
        text = re.sub(r'^#{1,6}\s+([0-9Xx-]{10,17}|TERMS OF USE)\s*$', r'\1', text, flags=re.MULTILINE)

        cleaned_lines = []
        for line in text.splitlines():
            match = re.match(r'^(#{1,6})\s+(.*)$', line)
            if match:
                hashes, raw_heading = match.groups()
                clean_text = clean_heading_text(raw_heading)
                if clean_text is not None:
                    cleaned_lines.append(f"{hashes} {clean_text}")
            else:
                cleaned_lines.append(line)

        lines = cleaned_lines
        first_heading_idx, first_heading_level, first_heading_text = None, None, None
        for idx, line in enumerate(lines):
            match = re.match(r'^(#{1,6})\s+(.*)$', line)
            if match:
                first_heading_idx = idx
                first_heading_level = len(match.group(1))
                first_heading_text = match.group(2).strip()
                break

        has_top_l1 = False
        if first_heading_idx is not None and first_heading_level == 1:
            if book_title:
                t_norm = book_title.lower().strip()
                h_norm = first_heading_text.lower().strip()
                if h_norm == t_norm or h_norm in t_norm or t_norm in h_norm:
                    has_top_l1 = True
            else:
                if first_heading_idx == 0 or all(not l.strip() for l in lines[:first_heading_idx]):
                    has_top_l1 = True

        stack = [(1, 1)] if (has_top_l1 or book_title) else [(0, 0)]
        final_lines = []
        if not has_top_l1 and book_title:
            final_lines.append(f"# {book_title}\n")

        for idx, line in enumerate(lines):
            if has_top_l1 and idx == first_heading_idx:
                final_lines.append(f"# {first_heading_text}")
                continue

            match = re.match(r'^(#{1,6})\s+(.*)$', line)
            if not match:
                final_lines.append(line)
                continue

            hashes, heading_content = match.groups()
            raw_lvl = len(hashes)
            if re.match(r'^Chapter\s+([0-9]+|[IVXLCDM]+):', heading_content, re.IGNORECASE):
                raw_lvl = 2

            if raw_lvl > stack[-1][0]:
                new_assigned = min(6, stack[-1][1] + 1)
                stack.append((raw_lvl, new_assigned))
            elif raw_lvl == stack[-1][0]:
                new_assigned = stack[-1][1]
            else:
                while len(stack) > 1 and stack[-1][0] > raw_lvl:
                    stack.pop()
                if stack[-1][0] == raw_lvl:
                    new_assigned = stack[-1][1]
                else:
                    new_assigned = min(6, stack[-1][1] + 1)
                    stack.append((raw_lvl, new_assigned))

            final_lines.append(f"{'#' * new_assigned} {heading_content.strip()}")

        result = "\n".join(final_lines).strip()
        return re.sub(r'\n{3,}', '\n\n', result)

    stem = re.sub(r'^[a-z0-9.]+[_-]', '', path.stem, flags=re.IGNORECASE)
    book_title = stem.replace("-", " ").replace("_", " ").title()

    return normalize_markdown_headings(raw_md, book_title=book_title)


def _convert_docx(path: Path, config=None, errors_out: Optional[List[dict]] = None) -> str:
    try:
        import docx
    except ImportError:
        raise ImportError("python-docx is required for converting .docx files. Install with `pip install python-docx`.")

    doc = docx.Document(str(path))
    md_lines: List[str] = []

    for elem in doc.element.body:
        tag = elem.tag.split('}')[-1]
        if tag == 'p':
            p = docx.text.paragraph.Paragraph(elem, doc)
            
            # Extract images embedded within the paragraph using Vision API or Multimodal LLM if configured
            if config and _is_image_processing_enabled(config):
                p_images = []
                for blip in elem.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}blip'):
                    embed_id = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if embed_id and embed_id in doc.part.related_parts:
                        image_part = doc.part.related_parts[embed_id]
                        image_bytes = image_part.blob
                        filename = getattr(image_part, "filename", "image.png")
                        p_images.append((image_bytes, filename))
                if p_images:
                    results = _process_images_concurrently(p_images, config, errors_out=errors_out)
                    for vision_md in results:
                        if vision_md:
                            md_lines.append(vision_md + "\n")
            
            text = _extract_text_and_math(elem).strip()
            if not text:
                continue
            style_name = p.style.name.lower() if p.style else ""
            if "heading 1" in style_name:
                md_lines.append(f"# {text}\n")
            elif "heading 2" in style_name:
                md_lines.append(f"## {text}\n")
            elif "heading 3" in style_name:
                md_lines.append(f"### {text}\n")
            elif "heading 4" in style_name:
                md_lines.append(f"#### {text}\n")
            elif "heading" in style_name:
                md_lines.append(f"##### {text}\n")
            elif "list" in style_name or "bullet" in style_name:
                md_lines.append(f"- {text}")
            else:
                md_lines.append(f"{text}\n")

        elif tag == 'tbl':
            tbl = docx.table.Table(elem, doc)
            table_md = _format_matrix_to_md_table([
                [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                for row in tbl.rows
            ])
            if table_md:
                md_lines.append(table_md + "\n")

    return "\n".join(md_lines).strip()


def _convert_pptx(path: Path, config=None, errors_out: Optional[List[dict]] = None) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required for converting .pptx files. Install with `pip install python-pptx`.")

    prs = Presentation(str(path))
    md_lines: List[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_title = ""
        slide_texts = []
        slide_images = []

        for shape in slide.shapes:
            if hasattr(shape, "image") and config and _is_image_processing_enabled(config):
                try:
                    image_bytes = shape.image.blob
                    filename = getattr(shape.image, "filename", "slide_image.png")
                    slide_images.append((image_bytes, filename))
                except Exception:
                    pass
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if shape == slide.shapes.title:
                    if hasattr(shape, "element"):
                        parsed_title = _extract_text_and_math(shape.element).strip()
                        slide_title = parsed_title.replace('\n', ' ')
                    else:
                        slide_title = text.replace('\n', ' ')
                else:
                    for paragraph in shape.text_frame.paragraphs:
                        if hasattr(paragraph, "_p"):
                            p_text = _extract_text_and_math(paragraph._p).strip()
                        else:
                            p_text = paragraph.text.strip()
                        if p_text:
                            level = paragraph.level
                            indent = "  " * level
                            slide_texts.append(f"{indent}- {p_text}")
            elif shape.has_table:
                table_matrix = []
                for row in shape.table.rows:
                    table_matrix.append([cell.text.strip().replace('\n', ' ') for cell in row.cells])
                table_md = _format_matrix_to_md_table(table_matrix)
                if table_md:
                    slide_texts.append(table_md)

        if slide_images:
            results = _process_images_concurrently(slide_images, config, errors_out=errors_out)
            for vision_md in results:
                if vision_md:
                    slide_texts.append(vision_md)

        header = f"## Slide {i}"
        if slide_title:
            header += f": {slide_title}"
        md_lines.append(header + "\n")

        if slide_texts:
            md_lines.extend(slide_texts)
            md_lines.append("")

    return "\n".join(md_lines).strip()


def _convert_xlsx(path: Path, config=None, errors_out: Optional[List[dict]] = None) -> str:
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl is required for converting .xlsx files. Install with `pip install openpyxl`.")

    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    md_lines: List[str] = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        md_lines.append(f"## Sheet: {sheet_name}\n")

        if config and _is_image_processing_enabled(config):
            sheet_images = []
            for img in getattr(sheet, "_images", []):
                try:
                    if hasattr(img, 'ref') and hasattr(img.ref, 'getvalue'):
                        image_bytes = img.ref.getvalue()
                    elif hasattr(img, '_data') and callable(img._data):
                        image_bytes = img._data()
                    else:
                        continue
                    filename = "excel_image.png"
                    sheet_images.append((image_bytes, filename))
                except Exception:
                    pass
            if sheet_images:
                results = _process_images_concurrently(sheet_images, config, errors_out=errors_out)
                for vision_md in results:
                    if vision_md:
                        md_lines.append(vision_md + "\n")

        matrix = []
        for row in sheet.iter_rows(values_only=True):
            if not row or all(v is None for v in row):
                continue
            row_vals = [str(v).strip().replace('\n', ' ') if v is not None else "" for v in row]
            matrix.append(row_vals)

        table_md = _format_matrix_to_md_table(matrix)
        if table_md:
            md_lines.append(table_md + "\n")

    wb.close()
    return "\n".join(md_lines).strip()


def _convert_csv(path: Path) -> str:
    try:
        content = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        content = path.read_text(encoding="latin-1", errors="replace")

    reader = csv.reader(io.StringIO(content))
    matrix = []
    for row in reader:
        if row and any(cell.strip() for cell in row):
            matrix.append([cell.strip().replace('\n', ' ') for cell in row])

    return _format_matrix_to_md_table(matrix)


def _convert_epub(path: Path) -> str:
    from qmd.epub import convert_epub_to_markdown
    try:
        return convert_epub_to_markdown(path)
    except Exception as e:
        raise ValueError(f"Failed to parse EPUB file '{path.name}': {e}")


def _convert_html(path: Path) -> str:
    try:
        content = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        content = path.read_text(encoding="latin-1", errors="replace")

    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("beautifulsoup4 is required for converting .html files. Install with `pip install beautifulsoup4`.")

    soup = BeautifulSoup(content, 'html.parser')

    for element in soup(["script", "style", "head", "title", "meta", "noscript", "svg", "iframe"]):
        element.decompose()

    md_lines: List[str] = []

    def process_node(node):
        if isinstance(node, str):
            text = node.strip()
            if text:
                md_lines.append(text)
            return

        name = node.name.lower() if node.name else ""

        if name in ["h1", "h2", "h3", "h4", "h5", "h6"]:
            level = int(name[1])
            prefix = "#" * level
            md_lines.append(f"\n{prefix} {node.get_text(strip=True)}\n")
        elif name == "p":
            text = node.get_text(strip=True)
            if text:
                md_lines.append(f"\n{text}\n")
        elif name in ["ul", "ol"]:
            for li in node.find_all("li", recursive=False):
                li_text = li.get_text(strip=True)
                if li_text:
                    md_lines.append(f"- {li_text}")
            md_lines.append("")
        elif name == "table":
            matrix = []
            for tr in node.find_all("tr"):
                row = [td.get_text(strip=True).replace('\n', ' ') for td in tr.find_all(["th", "td"])]
                if row:
                    matrix.append(row)
            table_md = _format_matrix_to_md_table(matrix)
            if table_md:
                md_lines.append(f"\n{table_md}\n")
        elif name in ["pre", "code"]:
            code_text = node.get_text()
            md_lines.append(f"\n```\n{code_text}\n```\n")
        else:
            for child in node.children:
                process_node(child)

    process_node(soup.body or soup)

    result = "\n".join(md_lines)
    result = re.sub(r'\n{3,}', '\n\n', result)
    return result.strip()


def guess_document_date(file_path: Union[str, Path], markdown_body: str = "") -> Optional[str]:
    """
    Attempts to infer the document date from the filename/path and markdown content using dplib.
    """
    dplib_mod = _get_dplib()
    if dplib_mod is None:
        return None
    try:
        sample_content = markdown_body[:4000] if markdown_body else ""
        if hasattr(dplib_mod, "extract_date"):
            try:
                res = dplib_mod.extract_date(path=str(file_path), content=sample_content)
            except TypeError:
                res = dplib_mod.extract_date(str(file_path), sample_content)
            if res:
                if hasattr(res, "hour") and res.hour == 0 and res.minute == 0 and res.second == 0 and res.microsecond == 0:
                    return res.strftime("%Y-%m-%d")
                return res.isoformat() if hasattr(res, "isoformat") else str(res)
        elif hasattr(dplib_mod, "DateResolver"):
            resolver = dplib_mod.DateResolver()
            report = resolver.resolve(filename=str(file_path).replace("\\", "/"), text_content=sample_content)
            if report and report.resolved_date:
                dt = report.resolved_date
                if dt.hour == 0 and dt.minute == 0 and dt.second == 0 and dt.microsecond == 0:
                    return dt.strftime("%Y-%m-%d")
                return dt.isoformat()
    except Exception:
        pass
    return None


def _format_matrix_to_md_table(matrix: List[List[str]]) -> str:
    if not matrix:
        return ""

    max_cols = max(len(row) for row in matrix)
    if max_cols == 0:
        return ""

    norm_matrix = []
    for row in matrix:
        norm_row = row + [""] * (max_cols - len(row))
        # Condense multiple spaces/newlines into a single space, and escape pipes
        norm_row = [re.sub(r'\s+', ' ', cell).strip().replace("|", "\\|") for cell in norm_row]
        norm_matrix.append(norm_row)

    header = norm_matrix[0]
    separator = ["---"] * max_cols
    rows = norm_matrix[1:] if len(norm_matrix) > 1 else []

    header_line = "| " + " | ".join(header) + " |"
    sep_line = "| " + " | ".join(separator) + " |"

    body_lines = ["| " + " | ".join(row) + " |" for row in rows]

    return "\n".join([header_line, sep_line] + body_lines)


def main():
    import argparse
    import sys

    # Ensure parent package directory (src) is in sys.path if running standalone script
    src_dir = str(Path(__file__).resolve().parent.parent)
    if src_dir not in sys.path:
        sys.path.insert(0, src_dir)

    parser = argparse.ArgumentParser(
        description="Convert a document (.docx, .pptx, .xlsx, .csv, .html, .epub, .md) to Markdown and inspect parsed blocks/chunks."
    )
    parser.add_argument("file_path", type=str, help="Path to the document to convert")
    parser.add_argument("-o", "--output", type=str, help="Optional output path to save the Markdown content")
    parser.add_argument("--show-blocks", action="store_true", help="Display parsed semantic blocks from docparse")
    parser.add_argument("--show-chunks", action="store_true", help="Display chunked content ready for embedding")
    parser.add_argument("--vision-url", type=str, help="URL for the Vision API to test image extraction")
    parser.add_argument("--vision-api-key", type=str, help="Optional API key for the Vision API")
    parser.add_argument("--multimodal-url", type=str, help="URL for OpenAI-compatible multimodal endpoint")
    parser.add_argument("--multimodal-api-key", type=str, help="Optional API key for multimodal endpoint")
    parser.add_argument("--multimodal-model", type=str, help="Model name for multimodal endpoint")
    parser.add_argument("--max-image-concurrency", type=int, default=4, help="Max simultaneous images to process")

    args = parser.parse_args()

    file_path = Path(args.file_path).expanduser().resolve()
    if not file_path.exists():
        print(f"Error: File not found: {file_path}", file=sys.stderr)
        sys.exit(1)

    if not is_supported_file(file_path):
        print(f"Warning: Extension '{file_path.suffix}' is not explicitly supported. Attempting plain text conversion...", file=sys.stderr)

    mock_config = None
    if args.vision_url or args.multimodal_url or args.multimodal_model:
        class MockConfig:
            vision_url = args.vision_url
            vision_api_key = args.vision_api_key
            multimodal_url = args.multimodal_url
            multimodal_api_key = args.multimodal_api_key
            multimodal_model = args.multimodal_model
            max_image_concurrency = args.max_image_concurrency
            request_timeout = 120.0
        mock_config = MockConfig()

    try:
        md_content = convert_to_markdown(file_path, config=mock_config)
    except Exception as e:
        print(f"Error converting {file_path}: {e}", file=sys.stderr)
        sys.exit(1)

    print("=" * 80)
    print(f"CONVERTED MARKDOWN ({file_path.name})")
    print("=" * 80)
    print(md_content)
    print("=" * 80)
    print(f"Stats: {len(md_content)} characters, {len(md_content.splitlines())} lines")
    inferred_date = guess_document_date(file_path, md_content)
    date_str = inferred_date if inferred_date else "None detected"
    print(f"Inferred Date: {date_str}")

    if args.output:
        out_path = Path(args.output).expanduser().resolve()
        out_path.write_text(md_content, encoding="utf-8")
        print(f"\nSaved output to: {out_path}")

    if args.show_blocks or args.show_chunks:
        try:
            from qmd.docparse.parser import parse_markdown_to_blocks
            from qmd.docparse.grouper import group_blocks_into_chunks

            blocks, has_tables = parse_markdown_to_blocks(content=md_content)
            print("\n" + "=" * 80)
            print(f"PARSED SEMANTIC BLOCKS ({len(blocks)} blocks found, has_tables={has_tables})")
            print("=" * 80)
            for i, block in enumerate(blocks, 1):
                print(f" Block #{i} [H{block.header_level}: {block.header_text}] (line {block.start_line}):")
                preview = block.content[:150].replace('\n', ' ')
                print(f"   {preview}..." if len(block.content) > 150 else f"   {preview}")

            if args.show_chunks:
                chunks = group_blocks_into_chunks(blocks, max_chunk_size=2048, target_chunk_size=1024)
                print("\n" + "=" * 80)
                print(f"GENERATED CHUNKS ({len(chunks)} chunks generated)")
                print("=" * 80)
                for i, chunk in enumerate(chunks, 1):
                    headers = " > ".join(chunk.parent_headers.values()) if chunk.parent_headers else "Root"
                    print(f"\n--- Chunk #{i} (Context: {headers}) ---")
                    print(chunk.content)
        except Exception as e:
            print(f"\nError running docparse analysis: {e}", file=sys.stderr)


if __name__ == "__main__":
    main()